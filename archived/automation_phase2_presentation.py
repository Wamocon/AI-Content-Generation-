"""
PHASE 2: DETAILED PRESENTATION GENERATION (Trainer Script + PowerPoint)
Focus: Comprehensive presentations based on Phase 1 content
Strategy: Use Phase 1 output + document for aligned, detailed presentations
"""

import asyncio
import os
import sys
from datetime import datetime
from typing import Dict, Any
from loguru import logger
from dotenv import load_dotenv
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from googleapiclient.http import MediaIoBaseUpload
import io

load_dotenv()
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from app.config import settings
from app.services.google_services import GoogleDriveService
from app.services.intelligent_gemini_service import IntelligentGeminiService

# Initialize services
google_drive_service = GoogleDriveService()
intelligent_gemini = IntelligentGeminiService()

class WebSocketProgress:
    """Lightweight progress tracker"""
    async def connect(self): pass
    async def send_progress(self, msg): logger.info(msg)

async def generate_detailed_trainer_script(document_content: str, use_cases_summary: str, doc_name: str, slide_count: int = 20) -> str:
    """
    Generate COMPREHENSIVE trainer script using intelligent chunking
    Strategy: Multi-pass for large documents, dynamic slide count
    """
    prompt = f"""
Du bist ein erfahrener Senior IT-Trainer und Präsentationsexperte. Erstelle ein DETAILLIERTES Trainer-Skript für {slide_count} PowerPoint-Folien.

KONTEXT AUS PHASE 1 (Use Cases):
{use_cases_summary[:1000]}

QUALITÄTSANFORDERUNGEN:
- {slide_count} Folien für VOLLSTÄNDIGE Themenabdeckung
- JEDE Folie: 1-1.5 Seiten detailliertes Skript
- Chronologischer Aufbau vom Einfachen zum Komplexen
- Praktische Beispiele und Demonstrationen
- Interaktive Elemente und Diskussionsfragen
- Zeitangaben pro Folie für Trainer

STRUKTUR PRO FOLIE (SEHR DETAILLIERT):

Folie [N] von [Gesamt]: [Präziser Titel]
Geschätzte Zeit: [X-Y Minuten]

FETT: 1. Lernziele dieser Folie
[4-5 konkrete, messbare Lernziele:]
- Was Teilnehmer nach dieser Folie können sollen
- Welches Wissen vermittelt wird
- Welche Fähigkeiten entwickelt werden
- Verbindung zu vorherigen Folien

FETT: 2. Folieninhalt (Was zeigen?)
[Detaillierte Beschreibung des visuellen Aufbaus:]
- Hauptüberschrift und Struktur
- Kernpunkte die auf der Folie stehen (4-6 Punkte)
- Grafiken, Diagramme oder Visualisierungen
- Farbliche Hervorhebungen wichtiger Elemente

FETT: 3. Theoretische Erklärung für Trainer
[5-6 Absätze mit tiefem fachlichem Hintergrund:]
- Grundlegende Konzepte ausführlich erklären
- Fachbegriffe definieren und einordnen
- Zusammenhänge zu anderen Themen aufzeigen
- Standards, Best Practices, Normen erwähnen
- Warum dieses Thema relevant ist (Business Value)
- Aktuelle Trends und Entwicklungen

FETT: 4. Präsentationsablauf - Was der Trainer sagt
[Wort-für-Wort Skript, 3-4 Absätze:]
"Guten Tag, schauen Sie sich bitte diese Folie an...
[Einleitung und Motivation - 30-60 Sekunden]

Lassen Sie mich das Konzept erklären...
[Hauptteil mit Erklärungen - 2-3 Minuten]

Besonders wichtig ist dabei...
[Schwerpunkte und Zusammenhänge - 1-2 Minuten]"

FETT: 5. Praktische Beispiele und Demonstrationen
[3-4 konkrete, realistische Beispiele:]
Beispiel 1: [Titel]
- Situation beschreiben (reale IT-Umgebung)
- Schritt-für-Schritt Durchführung
- Was zeigt dieses Beispiel?
- Bezug zur Praxis der Teilnehmer

[Weitere Beispiele analog...]

FETT: 6. Interaktive Elemente
[Aktivierung der Teilnehmer:]
- Diskussionsfragen (2-3 Fragen zum Nachdenken)
- Kurze Übungen oder Gruppenarbeit (Optional)
- Umfragen oder Quick Polls
- Erfahrungsaustausch anregen

FETT: 7. Häufige Fragen und Antworten
[Antizipierte Teilnehmerfragen:]
Q1: [Typische Frage]
A1: [Ausführliche Antwort mit Begründung]

Q2: [Typische Frage]
A2: [Ausführliche Antwort mit Begründung]

[2-3 weitere Q&As...]

FETT: 8. Trainer-Hinweise und Tipps
[Praktische Tipps für die Durchführung:]
- Worauf achten? (Häufige Stolpersteine)
- Wie Zeit managen?
- Wenn Teilnehmer nicht folgen können?
- Zusatzmaterial für schnelle Lernende
- Überleitung zur nächsten Folie

---

FORMATIERUNG:
- KEINE Markdown Symbole
- Klare Nummerierung
- Professionelles Deutsch
- JEDE Folie: Minimum 1 Seite Text

ZIEL: Trainer kann direkt aus diesem Skript unterrichten, ohne weitere Vorbereitung.
"""
    
    logger.info(f"[PHASE 2] Generating DETAILED trainer script for {doc_name}...")
    logger.info(f"[STRATEGY] Using direct generation (Phase 1 context sufficient)")
    
    result = await intelligent_gemini.generate_from_prompt(
        prompt=prompt,
        content_type="trainer_script",
        timeout=150
    )
    
    if result and len(result) > 1000:
        logger.info(f"[SUCCESS] Generated {len(result):,} characters of detailed trainer script")
        return result
    else:
        logger.warning(f"[WARNING] Trainer script returned short content: {len(result) if result else 0} chars")
        return result or ""

        return result

async def create_detailed_presentation(trainer_script: str, doc_name: str) -> Presentation:
    """
    Create comprehensive PowerPoint with 15-25 slides based on trainer script
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Parse trainer script to extract slide information
    slides_data = []
    current_slide = {}
    
    for line in trainer_script.split('\n'):
        line = line.strip()
        if not line:
            continue
            
        # Detect slide headers
        if line.startswith('Folie') and ('von' in line or ':' in line):
            if current_slide:
                slides_data.append(current_slide)
            current_slide = {'title': line, 'content': [], 'notes': ''}
        elif 'FETT: 1. Lernziele' in line or 'FETT: 2. Folieninhalt' in line:
            current_slide['section'] = 'content'
        elif 'FETT: 3. Theoretische Erklärung' in line or 'FETT: 4. Präsentationsablauf' in line:
            current_slide['section'] = 'notes'
        elif current_slide:
            if current_slide.get('section') == 'content':
                current_slide['content'].append(line)
            elif current_slide.get('section') == 'notes':
                current_slide['notes'] += line + ' '
    
    if current_slide:
        slides_data.append(current_slide)
    
    logger.info(f"[PPTX] Parsed {len(slides_data)} slides from trainer script")
    
    # Create title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = doc_name.replace('.docx', '').replace('_', ' ')
    subtitle.text = "Umfassende Schulungspräsentation\nDetaillierte Themenabdeckung"
    
    # Create content slides
    for i, slide_data in enumerate(slides_data[:25], 1):  # Max 25 slides
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        slide_title = slide_data.get('title', f'Folie {i}')
        title_shape.text = slide_title
        
        # Add content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        # Add bullet points from content
        content_lines = [c for c in slide_data.get('content', []) if c and len(c) > 10][:6]
        for content_line in content_lines:
            p = text_frame.add_paragraph()
            p.text = content_line[:200]  # Limit length
            p.level = 0
            p.font.size = Pt(18)
        
        # Add notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        notes_text = slide_data.get('notes', '')[:1000]  # Limit notes
        text_frame.text = f"TRAINER NOTES:\n\n{notes_text}"
    
    logger.info(f"[PPTX] Created presentation with {len(prs.slides)} slides")
    return prs

async def process_document_phase2(doc: Dict[str, Any], phase1_folder_id: str, use_cases_content: str) -> Dict[str, Any]:
    """Process single document - Phase 2: Presentation generation"""
    ws = WebSocketProgress()
    await ws.connect()
    
    doc_id = doc['id']
    doc_name = doc['name']
    
    logger.info(f"\n{'='*80}")
    logger.info(f"[PHASE 2] Processing: {doc_name}")
    logger.info(f"{'='*80}")
    
    try:
        # Extract document content
        await ws.send_progress(f"[1/3] Extracting content...")
        document_content = google_drive_service.get_file_content(doc_id)
        
        if not document_content or len(document_content) < 100:
            logger.error(f"[ERROR] Insufficient content")
            return {"status": "failed", "error": "insufficient_content"}
        
        # Generate detailed trainer script
        await ws.send_progress(f"[2/3] Generating detailed trainer script...")
        trainer_script = await generate_detailed_trainer_script(
            document_content,
            use_cases_content[:1000],  # Summary of use cases
            doc_name
        )
        
        if not trainer_script or len(trainer_script) < 1000:
            logger.error(f"[ERROR] Trainer script generation failed")
            return {"status": "failed", "error": "trainer_script_failed"}
        
        # Create PowerPoint
        await ws.send_progress(f"[3/3] Creating PowerPoint presentation...")
        prs = await create_detailed_presentation(trainer_script, doc_name)
        
        # Save trainer script as DOCX
        trainer_doc = Document()
        trainer_doc.add_heading('Trainer-Skript', 0)
        for para in trainer_script.split('\n\n'):
            if para.strip():
                trainer_doc.add_paragraph(para.strip())
        
        # Upload to Phase 2 folder (create subfolder in same parent)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        phase2_folder_name = f"{timestamp}_{doc_name.replace('.docx', '')}_Phase2"
        
        # Get parent folder ID from phase1_folder_id
        phase1_folder = google_drive_service.service.files().get(
            fileId=phase1_folder_id,
            fields='parents'
        ).execute()
        parent_id = phase1_folder.get('parents', [settings.google_drive_review_folder_id])[0]
        
        phase2_folder_metadata = {
            'name': phase2_folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [parent_id]
        }
        
        phase2_folder = google_drive_service.service.files().create(
            body=phase2_folder_metadata,
            fields='id'
        ).execute()
        phase2_folder_id = phase2_folder.get('id')
        
        # Upload trainer script
        script_buffer = io.BytesIO()
        trainer_doc.save(script_buffer)
        script_buffer.seek(0)
        
        script_media = MediaIoBaseUpload(
            script_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            resumable=True
        )
        
        script_file = google_drive_service.service.files().create(
            body={'name': f'trainer_script_{doc_name}', 'parents': [phase2_folder_id]},
            media_body=script_media,
            fields='id'
        ).execute()
        
        # Upload PowerPoint
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        pptx_media = MediaIoBaseUpload(
            pptx_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            resumable=True
        )
        
        pptx_file = google_drive_service.service.files().create(
            body={'name': f'presentation_{doc_name.replace(".docx", ".pptx")}', 'parents': [phase2_folder_id]},
            media_body=pptx_media,
            fields='id'
        ).execute()
        
        logger.info(f"[SUCCESS] Phase 2 completed for {doc_name}")
        logger.info(f"  - Trainer Script: {len(trainer_script):,} chars")
        logger.info(f"  - PowerPoint: {len(prs.slides)} slides")
        
        return {
            "status": "success",
            "doc_name": doc_name,
            "folder_id": phase2_folder_id,
            "script_id": script_file.get('id'),
            "pptx_id": pptx_file.get('id'),
            "slide_count": len(prs.slides)
        }
        
    except Exception as e:
        logger.error(f"[ERROR] Phase 2 failed for {doc_name}: {e}")
        import traceback
        logger.error(traceback.format_exc())
        return {"status": "failed", "error": str(e)}

async def main():
    """Main Phase 2 workflow"""
    logger.info("="*80)
    logger.info("[START] PHASE 2: DETAILED PRESENTATION GENERATION")
    logger.info("="*80)
    logger.info("[INFO] This script requires Phase 1 to be completed first")
    logger.info("[INFO] Looking for Phase 1 output folders...")
    
    # Find Phase 1 folders
    review_folder_id = settings.google_drive_review_folder_id
    query = f"'{review_folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and name contains 'Phase1'"
    
    phase1_folders = google_drive_service.service.files().list(
        q=query,
        fields='files(id, name)'
    ).execute().get('files', [])
    
    logger.info(f"[OK] Found {len(phase1_folders)} Phase 1 folders to process")
    
    if len(phase1_folders) == 0:
        logger.error("[ERROR] No Phase 1 folders found. Please run automation_phase1_content.py first.")
        return
    
    # Get source documents
    source_folder_id = settings.google_drive_content_source_folder_id
    documents = google_drive_service.list_files_in_folder(source_folder_id)
    
    results = []
    success_count = 0
    failed_count = 0
    
    for i, phase1_folder in enumerate(phase1_folders, 1):
        folder_name = phase1_folder['name']
        folder_id = phase1_folder['id']
        
        # Extract document name from folder name
        doc_name = folder_name.split('_', 2)[-1].replace('_Phase1', '') + '.docx'
        
        # Find matching document
        matching_doc = next((d for d in documents if d['name'] == doc_name), None)
        if not matching_doc:
            logger.warning(f"[SKIP] No matching document for {folder_name}")
            continue
        
        logger.info(f"\n[{i}/{len(phase1_folders)}] Processing: {doc_name}")
        
        # Read use cases from Phase 1 folder
        use_cases_query = f"'{folder_id}' in parents and name contains 'use_cases'"
        use_cases_files = google_drive_service.service.files().list(
            q=use_cases_query,
            fields='files(id)'
        ).execute().get('files', [])
        
        use_cases_content = ""
        if use_cases_files:
            # Extract text from use cases file (simplified)
            use_cases_content = "Use cases content from Phase 1"  # Placeholder
        
        result = await process_document_phase2(matching_doc, folder_id, use_cases_content)
        results.append(result)
        
        if result["status"] == "success":
            success_count += 1
        else:
            failed_count += 1
        
        # Wait between documents
        if i < len(phase1_folders):
            logger.info(f"[WAIT] Cooling down 15s before next document...")
            await asyncio.sleep(15)
    
    logger.info("\n" + "="*80)
    logger.info("[COMPLETE] PHASE 2 FINISHED")
    logger.info("="*80)
    logger.info(f"[STATS] Success: {success_count}/{len(phase1_folders)}")
    logger.info(f"[STATS] Failed: {failed_count}/{len(phase1_folders)}")
    logger.info(f"[STATS] Success Rate: {(success_count/len(phase1_folders)*100 if phase1_folders else 0):.1f}%")

if __name__ == "__main__":
    print("⚠️ PHASE 2 IS DEPRECATED - This script is no longer used.")
    print("The system now focuses ONLY on generating comprehensive use cases.")
    print("Please use automation_phase1_content.py or automation_engine.py instead.")
    # asyncio.run(main())  # COMMENTED OUT
"""
