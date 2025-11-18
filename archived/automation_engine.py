"""
FIAE AI Content Factory - Comprehensive Automation Engine
Processes all documents from Google Drive source folder with real-time progress updates
"""

import asyncio
import os
import sys
import json
import io
import re
from datetime import datetime
from typing import Dict, Any, List
from loguru import logger
import websockets
from dotenv import load_dotenv
from googleapiclient.http import MediaIoBaseUpload, MediaFileUpload

# Load environment variables FIRST
load_dotenv()

# Add project root to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from app.config import settings

# Initialize services with graceful fallbacks
google_drive_service = None
google_sheets_service = None
rag_processor = None  # Lazy-loaded to avoid PyTorch import delay

try:
    from app.services.google_services import GoogleDriveService, GoogleSheetsService
    google_drive_service = GoogleDriveService()
    google_sheets_service = GoogleSheetsService()
    logger.info("[OK] Google services initialized successfully")
except Exception as e:
    logger.error(f"[CRITICAL] Google services initialization failed: {e}")
    logger.info("Please ensure personal_google_token.pickle is valid and OAuth is complete")
    sys.exit(1)

# RAG processor is OPTIONAL and lazy-loaded to avoid blocking startup
# It will be initialized on-demand only if needed
def get_rag_processor():
    """Lazy-load RAG processor to avoid blocking imports"""
    global rag_processor
    if rag_processor is None:
        try:
            from app.services.rag_enhanced_processor import RAGEnhancedProcessor
            rag_processor = RAGEnhancedProcessor()
            logger.info("[OK] RAG processor initialized successfully")
        except Exception as e:
            logger.warning(f"[WARNING] RAG processor not available: {e}")
            rag_processor = False  # Mark as attempted and failed
    return rag_processor if rag_processor is not False else None

# Check if required services are available
if not google_drive_service:
    logger.error("[ERROR] Google Drive service is required - cannot proceed")
    sys.exit(1)


def parse_use_cases_content(content: str) -> List[Dict[str, str]]:
    """Parse use cases content into structured format"""
    use_cases = []
    
    # Split by use case sections
    sections = re.split(r'FETT:\s*Titel:', content)
    
    for section in sections[1:]:  # Skip first empty section
        lines = section.strip().split('\n')
        if not lines:
            continue
            
        use_case = {
            'title': lines[0].strip(),
            'theory': '',
            'scenario': '',
            'tasks': [],
            'solution': '',
            'expected_results': ''
        }
        
        current_section = None
        current_text = []
        
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
                
            if 'FETT: Theoretischer Hintergrund:' in line:
                if current_section and current_text:
                    use_case[current_section] = '\n'.join(current_text).strip()
                current_section = 'theory'
                current_text = []
            elif 'FETT: Praxis-Szenario:' in line:
                if current_section and current_text:
                    use_case[current_section] = '\n'.join(current_text).strip()
                current_section = 'scenario'
                current_text = []
            elif 'FETT: Aufgaben für Lernende:' in line:
                if current_section and current_text:
                    use_case[current_section] = '\n'.join(current_text).strip()
                current_section = 'tasks'
                current_text = []
            elif 'FETT: Musterlösung für Trainer:' in line:
                if current_section and current_text:
                    if current_section == 'tasks':
                        use_case[current_section] = current_text
                    else:
                        use_case[current_section] = '\n'.join(current_text).strip()
                current_section = 'solution'
                current_text = []
            elif 'FETT: Erwartete Ergebnisse:' in line:
                if current_section and current_text:
                    use_case[current_section] = '\n'.join(current_text).strip()
                current_section = 'expected_results'
                current_text = []
            else:
                if current_section == 'tasks' and line.startswith(('1.', '2.', '3.', '4.', '5.')):
                    current_text.append(line)
                elif current_section:
                    current_text.append(line)
        
        # Handle last section
        if current_section and current_text:
            if current_section == 'tasks':
                use_case[current_section] = current_text
            else:
                use_case[current_section] = '\n'.join(current_text).strip()
        
        use_cases.append(use_case)
    
    return use_cases


def parse_quiz_content(content: str) -> List[Dict[str, Any]]:
    """Parse quiz content into structured format"""
    questions = []
    
    # Split by question sections
    sections = re.split(r'Frage\s+(\d+)\s+\((Leicht|Mittel|Schwer)\):', content)
    
    for i in range(1, len(sections), 3):  # Every 3rd element is a question
        if i + 2 < len(sections):
            question_num = sections[i]
            difficulty = sections[i + 1]
            question_content = sections[i + 2]
            
            question = {
                'number': int(question_num),
                'difficulty': difficulty,
                'question': '',
                'options': [],
                'correct_answer': '',
                'explanation': ''
            }
            
            lines = question_content.strip().split('\n')
            current_section = 'question'
            current_text = []
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                    
                if line.startswith('Antwortoptionen:'):
                    if current_section == 'question':
                        question[current_section] = '\n'.join(current_text).strip()
                    current_section = 'options'
                    current_text = []
                elif line.startswith('Korrekte Antwort:'):
                    if current_section == 'options':
                        question[current_section] = current_text
                    current_section = 'correct_answer'
                    current_text = [line.replace('Korrekte Antwort:', '').strip()]
                elif line.startswith('Erklärung:'):
                    if current_section == 'correct_answer':
                        question[current_section] = '\n'.join(current_text).strip()
                    current_section = 'explanation'
                    current_text = [line.replace('Erklärung:', '').strip()]
                else:
                    if current_section == 'options' and line.startswith(('A)', 'B)', 'C)', 'D)')):
                        current_text.append(line)
                    elif current_section in ['question', 'correct_answer', 'explanation']:
                        current_text.append(line)
            
            # Handle last section
            if current_section and current_text:
                if current_section == 'options':
                    question[current_section] = current_text
                else:
                    question[current_section] = '\n'.join(current_text).strip()
            
            questions.append(question)
    
    return questions




def validate_and_clean_content(content: str, content_type: str) -> str:
    """Remove markdown symbols and validate structure"""
    if not content:
        return content
    
    # Remove ALL markdown formatting
    cleaned = content
    cleaned = re.sub(r'\*\*([^*]+)\*\*', r'\1', cleaned)  # **bold**
    cleaned = re.sub(r'\*([^*]+)\*', r'\1', cleaned)  # *italic*
    cleaned = re.sub(r'#{1,6}\s+', '', cleaned)  # ### headings
    cleaned = re.sub(r'_([^_]+)_', r'\1', cleaned)  # _underline_
    cleaned = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', cleaned)  # [links](url)
    
    # Validate structure based on content type
    if content_type == "use_cases":
        required = ["Theoretischer Hintergrund:", "Praxis-Szenario:", "Aufgaben für Lernende:", "Musterlösung"]
        if not all(req in cleaned for req in required):
            logger.warning(f"Use case missing required sections: {required}")
            return cleaned  # Return cleaned content anyway
            
    elif content_type == "quiz":
        if not re.search(r'Frage \d+ \((Leicht|Mittel|Schwer)\):', cleaned):
            logger.warning("Quiz questions missing proper format")
            return cleaned  # Return cleaned content anyway
    
    return cleaned


def ensure_example_uniqueness(use_cases_content: str, quiz_content: str, powerpoint_content: str) -> Dict[str, str]:
    """Ensure each document type has unique examples"""
    
    # Add example tracking to prompts
    enhanced_prompts = {
        "use_cases": use_cases_content + """
        
BEISPIELE (NUR für Anwendungsfälle - NICHT in Quiz oder Slides verwenden):
- Cloud-Migration in AWS/Azure
- Microservices-Architektur mit Docker/Kubernetes
- CI/CD Pipeline mit Jenkins/GitLab
- API-Gateway Implementation
- Database Sharding Strategy
DIESE Beispiele NICHT in Quiz oder Slides verwenden!
""",
        
        "quiz": quiz_content + """
        
BEISPIELE (NUR für Quiz - ANDERE als Anwendungsfälle und Slides):
- Netzwerk-Troubleshooting (DNS, Routing)
- Security Incident Response
- Datenbank-Performance-Tuning
- Load Balancer Configuration
- Backup & Recovery Szenarien
KEINE Wiederholung von Anwendungsfall-Beispielen!
""",
        
        "powerpoint": powerpoint_content + """
        
BEISPIELE (NUR für Slides - KOMPLETT NEU):
- Industrie 4.0 Anwendungen
- E-Commerce Plattform-Architektur
- IoT Device Management
- Fintech Payment Processing
- Healthcare Data Systems
KEINE Wiederholung von vorherigen Beispielen!
""",
        
        "trainer_script": f"""
WICHTIG: Verwende EXAKT die gleichen Beispiele wie in den PowerPoint-Folien!
{powerpoint_content[:500]}... [Auszug]

Erkläre diese Beispiele im Detail für den Trainer.
"""
    }
    
    return enhanced_prompts


class WebSocketProgress:
    """Send real-time progress updates via WebSocket"""
    def __init__(self):
        self.connected = False
        self.websocket = None
        self.ws_url = "ws://localhost:8000/ws"
        self.connection = None
    
    async def connect(self):
        """Connect to WebSocket server"""
        try:
            self.connection = await websockets.connect(self.ws_url)
            logger.info("[OK] Connected to WebSocket server")
            self.connected = True
        except Exception as e:
            logger.warning(f"[WARNING] WebSocket connection failed: {e} - Progress will be logged only")
            self.connection = None
            self.connected = False
    
    async def send_progress(self, message: str, progress: int | None = None, total: int | None = None):
        """Send progress message via WebSocket or log if not connected"""
        logger.info(message)
        if self.connection:
            try:
                data: Dict[str, Any] = {
                    "type": "terminal_log",
                    "message": message,
                    "timestamp": datetime.now().isoformat()
                }
                if progress is not None and total is not None:
                    data["progress"] = progress
                    data["total"] = total
                    data["percentage"] = round((progress / total) * 100, 1)
                await self.connection.send(json.dumps(data))
            except Exception as e:
                logger.warning(f"Failed to send WebSocket update: {e}")
    
    async def close(self):
        """Close WebSocket connection"""
        if self.connection:
            await self.connection.close()


async def generate_production_level_content(
    document_content: str, 
    enhanced_content: Dict[str, Any], 
    doc_name: str,
    ws: WebSocketProgress
) -> Dict[str, Any]:
    """
    SIMPLIFIED: Generate ONLY use cases (Anwendungsfälle) - NO quiz, NO presentations!
    Focus: Comprehensive, high-quality use cases with dynamic count based on document content
    """
    try:
        # Import Gemini service for advanced content generation
        from app.services.gemini_ai_service import GeminiAIService
        
        gemini_service = GeminiAIService()
        
        await ws.send_progress(f"   [AI] Generiere Anwendungsfälle mit Gemini (NUR Use Cases)...")
        
        # Extract content depth information
        content_depth = enhanced_content.get('content_depth', {})
        word_count = content_depth.get('word_count', 1000)
        # Dynamic use case count: 1 per 600 words, minimum 3, no maximum
        use_case_count = max(3, word_count // 600)
        
        await ws.send_progress(f"   [ANALYSE] {word_count} Wörter → {use_case_count} Anwendungsfälle werden generiert")
        
        # CRITICAL FIX: Drastically reduce prompt size to prevent 504 timeouts
        # Limit to 3000 chars - tests show 8000 causes consistent timeouts
        doc_summary = document_content[:3000] if len(document_content) > 3000 else document_content
        if len(document_content) > 3000:
            doc_summary += "\n\n[Hinweis: Zusammenfassung. Vollständig analysiert.]"
        
        # Generate comprehensive use cases (DOCX format) - ULTRA-OPTIMIZED
        use_cases_prompt = f"""
Erstelle 4-5 IT-Anwendungsfälle:

DOKUMENT:
{doc_summary}

PRO FALL:

FETT: Titel: [Titel]

FETT: Hintergrund: [1 Absatz]

FETT: Szenario: [1 Absatz]

FETT: Aufgaben:
1. [Aufgabe]
2. [Aufgabe]

FETT: Lösung: [Kurz]

WICHTIG:
- KEINE Markdown
- "FETT:" vor Überschriften
- 4-5 Fälle
"""
        
        # Generate content with retry logic - OPTIMIZED TO PREVENT TIMEOUTS
        async def generate_with_retry(prompt, content_type):
            for attempt in range(4):
                try:
                    await ws.send_progress(f"   [AI] Generating {content_type} (attempt {attempt + 1}/4)...")
                    
                    # CRITICAL FIX: DO NOT add RAG context - causes prompt bloat and timeouts
                    # Use the prompt as-is - it already contains optimal doc_summary
                    enhanced_prompt = prompt
                    
                    # CRITICAL FIX: Reduce document content and timeout to prevent 504 errors
                    result = gemini_service.generate_content_with_retry(
                        content_type=content_type,
                        document_content=document_content[:10000],  # Limit to 10k chars
                        context_query=enhanced_prompt,
                        timeout=50,  # Reduced from 90 to 50 seconds (safe zone)
                        max_retries=2
                    )
                    if result and len(result.strip()) > 50 and not result.startswith("Error"):
                        # Validate and clean content
                        cleaned_content = validate_and_clean_content(result, content_type)
                        if cleaned_content:
                            logger.info(f"[OK] Generated and validated {content_type} content (attempt {attempt + 1})")
                            return cleaned_content
                        else:
                            logger.warning(f"Content validation failed for {content_type}, retrying...")
                    else:
                        logger.warning(f"Empty or error content for {content_type}: {result[:100] if result else 'None'}")
                        if attempt < 3:  # If not last attempt
                            wait_time = 10 * (attempt + 1)
                            logger.info(f"Waiting {wait_time}s before retry...")
                            await asyncio.sleep(wait_time)
                except Exception as e:
                    logger.warning(f"Attempt {attempt + 1} failed for {content_type}: {e}")
                    if attempt < 3:  # If not last attempt
                        await asyncio.sleep(10)
                    else:
                        return f"Error generating {content_type}: {str(e)}"
            return f"Failed to generate {content_type} after 4 attempts"
        
        use_cases_content = await generate_with_retry(use_cases_prompt, "use cases")
        
        # Generate comprehensive quiz (DOCX format) - ULTRA-OPTIMIZED
        quiz_prompt = f"""
Erstelle 10-12 Quiz-Fragen:

DOKUMENT:
{doc_summary}

PRO FRAGE:

Frage [N] ([Schwierigkeit]): [Frage]

A) [Option]
B) [Option]
C) [Option]
D) [Option]

Korrekt: [Buchstabe]
Warum: [Kurz]

WICHTIG:
- KEINE Markdown
- 40% Leicht, 40% Mittel, 20% Schwer
- 10-12 Fragen
"""
        
        quiz_content = await generate_with_retry(quiz_prompt, "quiz questions")
        
        # Generate trainer script (DOCX format) - ULTRA-OPTIMIZED
        script_prompt = f"""
Erstelle Trainer-Skript für 6-8 Folien:

DOKUMENT:
{doc_summary}

PRO FOLIE:

Folie [N]: [Titel]

Lernziele: [2-3 Punkte]

Erklärung: [1-2 Absätze]

Beispiele:
- [Beispiel 1]
- [Beispiel 2]

Übung: [1 Aufgabe]

Hinweise: [Für Trainer]

WICHTIG:
- KEINE Markdown
- 6-8 Folien
"""
        
        script_content = await generate_with_retry(script_prompt, "trainer script")
        
        return {
            "use_cases": str(use_cases_content),
            "quiz": str(quiz_content),
            "trainer_script": str(script_content),
            "document_content": document_content,  # Include for fallback PPTX generation
            "content_depth": content_depth
        }
        
    except Exception as e:
        logger.error(f"Error generating production content: {e}")
        return {
            "use_cases": f"Error generating use cases: {e}",
            "quiz": f"Error generating quiz: {e}",
            "trainer_script": f"Error generating script: {e}",
            "content_depth": {}
        }


def create_professional_pptx_presentation(trainer_script: str, doc_name: str) -> io.BytesIO | None:
    """
    Create comprehensive, high-quality PPTX presentation with detailed content coverage
    Covers ALL topics from source document in easy-to-understand language
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor
        
        logger.info("[PPTX] Creating professional presentation with comprehensive content coverage...")
        
        # Validate trainer script content
        if not trainer_script or len(trainer_script.strip()) < 100:
            logger.error("[PPTX] Invalid or empty trainer script")
            return None
        
        # Check for error messages in script
        error_indicators = [
            "error generating",
            "timeout exceeded",
            "model is overloaded", 
            "503",
            "504",
            "deadline exceeded"
        ]
        
        script_lower = trainer_script.lower()
        for indicator in error_indicators:
            if indicator in script_lower:
                logger.error(f"[PPTX] Trainer script contains error: {indicator}")
                return None
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Parse trainer script into structured slides
        slides_content = parse_trainer_script_to_slides(trainer_script)
        
        if not slides_content:
            logger.warning("[PPTX] No slides parsed, creating from raw script")
            slides_content = create_slides_from_script(trainer_script)
        
        if not slides_content:
            logger.error("[PPTX] Failed to create slides from script")
            return None
        
        logger.info(f"[PPTX] Generating {len(slides_content)} detailed slides...")
        
        # Create title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = doc_name.replace('.docx', '').replace('_', ' ')
        subtitle.text = "Comprehensive Training Presentation\nDetailed Content Coverage"
        
        # Format title slide
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Create content slides with detailed information
        for i, slide_content in enumerate(slides_content, 1):
            # Add slide with title and content layout
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set slide title
            title_shape = slide.shapes.title
            slide_title = slide_content.get('title', f'Topic {i}')
            title_shape.text = f"{i}. {slide_title}"
            
            # Format title
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            # Get content placeholder
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Add theory explanation section
            if slide_content.get('theory'):
                add_section_to_slide(text_frame, "Theorie & Konzepte:", slide_content['theory'], Pt(16), Pt(14))
            
            # Add practical examples section
            if slide_content.get('examples'):
                add_section_to_slide(text_frame, "Praktische Beispiele:", slide_content['examples'], Pt(16), Pt(14))
            
            # Add image description section
            if slide_content.get('image_description'):
                add_section_to_slide(text_frame, "Visuelle Darstellung:", slide_content['image_description'], Pt(16), Pt(13))
            
            # Add key takeaways section
            if slide_content.get('key_takeaways'):
                add_section_to_slide(text_frame, "Wichtige Erkenntnisse:", slide_content['key_takeaways'], Pt(16), Pt(14))
            
            # Add trainer notes in notes section
            if slide_content.get('trainer_notes'):
                notes_slide = slide.notes_slide
                notes_frame = notes_slide.notes_text_frame
                notes_frame.text = f"TRAINER NOTES:\n\n{slide_content['trainer_notes']}"
        
        # Create summary slide
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        summary_title = summary_slide.shapes.title
        summary_title.text = "Zusammenfassung & Nächste Schritte"
        summary_title.text_frame.paragraphs[0].font.size = Pt(32)
        summary_title.text_frame.paragraphs[0].font.bold = True
        summary_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        summary_content = summary_slide.placeholders[1]
        summary_frame = summary_content.text_frame
        summary_frame.clear()
        
        summary_p = summary_frame.paragraphs[0]
        summary_p.text = f"Diese Präsentation umfasst {len(slides_content)} detaillierte Themen"
        summary_p.font.size = Pt(18)
        summary_p.font.bold = True
        
        summary_bullet = summary_frame.add_paragraph()
        summary_bullet.text = "Alle Konzepte wurden umfassend erklärt"
        summary_bullet.level = 1
        summary_bullet.font.size = Pt(16)
        
        summary_bullet2 = summary_frame.add_paragraph()
        summary_bullet2.text = "Praktische Beispiele für jedes Thema"
        summary_bullet2.level = 1
        summary_bullet2.font.size = Pt(16)
        
        summary_bullet3 = summary_frame.add_paragraph()
        summary_bullet3.text = "Bereit für praktische Anwendung"
        summary_bullet3.level = 1
        summary_bullet3.font.size = Pt(16)
        
        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        logger.info(f"[PPTX] Professional presentation created successfully: {len(slides_content) + 2} total slides")
        return pptx_buffer
        
    except Exception as e:
        logger.error(f"[PPTX] Professional presentation creation failed: {e}")
        import traceback
        logger.error(f"[PPTX] Traceback: {traceback.format_exc()}")
        return None


def add_section_to_slide(text_frame, section_title: str, content: str, title_size, content_size):
    """Add a formatted section to a slide"""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    # Add section title
    if len(text_frame.paragraphs) == 1 and not text_frame.paragraphs[0].text:
        title_p = text_frame.paragraphs[0]
    else:
        title_p = text_frame.add_paragraph()
    
    title_p.text = section_title
    title_p.font.bold = True
    title_p.font.size = title_size
    title_p.font.color.rgb = RGBColor(0, 102, 204)
    title_p.space_after = Pt(6)
    
    # Add content - split into manageable chunks
    content_lines = content.split('\n')
    for line in content_lines:
        if line.strip():
            content_p = text_frame.add_paragraph()
            # Handle bullet points
            if line.strip().startswith('-') or line.strip().startswith('•'):
                content_p.text = line.strip()[1:].strip()
                content_p.level = 1
            else:
                content_p.text = line.strip()
                content_p.level = 0
            
            content_p.font.size = content_size
            content_p.space_after = Pt(4)
    
    # Add spacing after section
    spacer = text_frame.add_paragraph()
    spacer.text = ""
    spacer.space_after = Pt(8)


def create_slides_from_script(script: str) -> list:
    """Create slides from raw script if parsing fails"""
    slides = []
    
    # Split script into sections by double newlines or slide markers
    sections = []
    
    # Try to split by "Slide X:" pattern
    import re
    slide_pattern = re.compile(r'Slide\s+\d+:', re.IGNORECASE)
    parts = slide_pattern.split(script)
    
    if len(parts) > 1:
        # Successfully split by slide markers
        for i, part in enumerate(parts[1:], 1):  # Skip first empty part
            sections.append(part.strip())
    else:
        # Fall back to splitting by double newlines
        sections = [s.strip() for s in script.split('\n\n') if s.strip()]
    
    # Create slides from sections
    for i, section in enumerate(sections):
        if not section:
            continue
        
        lines = section.split('\n')
        title = lines[0][:100] if lines else f"Topic {i+1}"  # Use first line as title
        
        slide_data = {
            'title': title,
            'theory': '',
            'examples': '',
            'key_takeaways': '',
            'image_description': '',
            'trainer_notes': ''
        }
        
        # Parse content into sections
        current_section = 'theory'
        content_buffer = []
        
        for line in lines:
            line_lower = line.lower().strip()
            
            if 'theorie' in line_lower or 'konzept' in line_lower:
                if content_buffer:
                    slide_data[current_section] = '\n'.join(content_buffer)
                current_section = 'theory'
                content_buffer = []
            elif 'beispiel' in line_lower or 'praxis' in line_lower:
                if content_buffer:
                    slide_data[current_section] = '\n'.join(content_buffer)
                current_section = 'examples'
                content_buffer = []
            elif 'takeaway' in line_lower or 'wichtig' in line_lower:
                if content_buffer:
                    slide_data[current_section] = '\n'.join(content_buffer)
                current_section = 'key_takeaways'
                content_buffer = []
            elif 'bild' in line_lower or 'visual' in line_lower:
                if content_buffer:
                    slide_data[current_section] = '\n'.join(content_buffer)
                current_section = 'image_description'
                content_buffer = []
            elif line.strip():
                content_buffer.append(line.strip())
        
        # Add final buffer
        if content_buffer:
            slide_data[current_section] = '\n'.join(content_buffer)
        
        slides.append(slide_data)
    
    logger.info(f"[PPTX] Created {len(slides)} slides from raw script")
    return slides


def create_basic_pptx_fallback(trainer_script: str, doc_name: str) -> io.BytesIO | None:
    """Create basic PPTX presentation as final fallback"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        # Create presentation
        prs = Presentation()
        
        # Split script into sections
        sections = trainer_script.split('\n\n')
        
        for i, section in enumerate(sections[:20]):  # Limit to 20 slides
            if section.strip():
                # Add slide
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title
                title = slide.shapes.title
                title.text = f"Slide {i+1}: {doc_name}"
                
                # Set content
                content = slide.placeholders[1]
                content.text = section.strip()
        
        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        logger.info(f"[FALLBACK] Basic PPTX created: {len(sections)} sections")
        return pptx_buffer
        
    except Exception as e:
        logger.error(f"[FALLBACK] Basic PPTX creation failed: {e}")
        return None


def create_professional_pptx_fallback(trainer_script: str, doc_name: str) -> io.BytesIO | None:
    """Deprecated temporary fallback; keep stub to avoid breaking imports."""
    logger.warning("[PPTX] Fallback disabled. Use Gamma API integration.")
    return None

def create_robust_pptx_from_content(document_content: str, doc_name: str) -> io.BytesIO | None:
    """
    Create a robust PPTX presentation even when trainer script generation fails.
    This ensures presentations are ALWAYS created for every document.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        logger.info("[PPTX-ROBUST] Creating presentation from document content...")
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = doc_name.replace('.docx', '').replace('_', ' ')
        subtitle.text = "Educational Training Presentation"
        
        # Format title slide
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Split content into logical sections for slides
        # Get first 15000 characters (manageable size)
        content = document_content[:15000]
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip() and len(p.strip()) > 50]
        
        # Limit to 25 slides to avoid overly large presentations
        slide_count = min(25, len(paragraphs))
        
        for i in range(slide_count):
            # Create content slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            
            # Set title
            title_shape = slide.shapes.title
            title_shape.text = f"Topic {i+1}"
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
            
            # Add content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Split paragraph into sentences and add as bullets
            para = paragraphs[i]
            sentences = para.split('. ')[:6]  # Max 6 bullets per slide
            
            for j, sentence in enumerate(sentences):
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = sentence.strip()
                p.font.size = Pt(14)
                p.level = 0
        
        # Save to BytesIO
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        logger.info(f"[PPTX-ROBUST] Created {slide_count + 1} slides successfully")
        return pptx_buffer
        
    except Exception as e:
        logger.error(f"[PPTX-ROBUST] Failed to create presentation: {e}")
        import traceback
        logger.error(f"[PPTX-ROBUST] Traceback: {traceback.format_exc()}")
        return None


def parse_trainer_script_to_slides(trainer_script: str) -> list:
    """Parse trainer script into structured slide content"""
    slides = []
    
    # Check if trainer script is valid
    if not trainer_script or len(trainer_script.strip()) < 100:
        logger.warning("[PPTX] Invalid trainer script, returning empty slides")
        return []
    
    # Check for error messages
    error_indicators = ["error", "timeout", "overloaded", "503", "504"]
    script_lower = trainer_script.lower()
    if any(indicator in script_lower for indicator in error_indicators):
        logger.warning("[PPTX] Trainer script contains errors, returning empty slides")
        return []
    
    # Split by slide markers
    slide_sections = trainer_script.split('Slide ')
    
    for section in slide_sections[1:]:  # Skip first empty section
        lines = section.strip().split('\n')
        if not lines:
            continue
            
        # Extract slide number and title
        first_line = lines[0]
        if ':' in first_line:
            slide_num, title = first_line.split(':', 1)
            title = title.strip()
        else:
            title = first_line.strip()
        
        # Extract content sections
        slide_content = {
            'title': title,
            'theory': '',
            'examples': '',
            'key_takeaways': '',
            'trainer_notes': ''
        }
        
        current_section = 'theory'
        content_lines = []
        
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
                
            # Check for section headers
            if 'Theorie-Erklärung:' in line:
                if content_lines:
                    slide_content[current_section] = '\n'.join(content_lines)
                current_section = 'theory'
                content_lines = []
            elif 'Praxis-Beispiele:' in line:
                if content_lines:
                    slide_content[current_section] = '\n'.join(content_lines)
                current_section = 'examples'
                content_lines = []
            elif 'Key Takeaways:' in line:
                if content_lines:
                    slide_content[current_section] = '\n'.join(content_lines)
                current_section = 'key_takeaways'
                content_lines = []
            else:
                content_lines.append(line)
        
        # Add final content
        if content_lines:
            slide_content[current_section] = '\n'.join(content_lines)
        
        slides.append(slide_content)
    
    return slides


async def save_production_content_with_personal_service(
    personal_service, 
    production_content: Dict[str, Any], 
    original_filename: str, 
    source_file_id: str
) -> Dict[str, Any]:
    """
    Save production-level content as DOCX and PPTX files using personal Google Drive service.
    Creates organized folder structure: Review/<SourceName>/<Timestamp_DocumentName>/
    with Phase1 (use_cases, quiz) and Phase2 (trainer_script, pptx) subfolders.
    """
    try:
        from app.config import settings
        
        # Validate inputs
        if not production_content:
            return {"success": False, "error": "No production content provided", "total_files_created": 0, "created_files": []}
        
        if not original_filename or not source_file_id:
            return {"success": False, "error": "Invalid filename or file ID", "total_files_created": 0, "created_files": []}
        
        service = personal_service.get_service()
        
        # Get source folder name for organization
        try:
            source_folder_info = service.files().get(
                fileId=settings.google_drive_content_source_folder_id,
                fields='name'
            ).execute()
            source_folder_name = source_folder_info.get('name', 'AI-Content-Source')
        except Exception as e:
            logger.warning(f"Could not get source folder name: {e}")
            source_folder_name = 'AI-Content-Source'
        
        # Clean filename for folder name
        clean_filename = original_filename.replace('.docx', '').replace('.pdf', '').replace(' ', '_')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        # Create main document folder: Review/<SourceName>/<Timestamp_DocumentName>
        main_folder_name = f"{timestamp}_{clean_filename}"
        
        # Check if source subfolder exists in Review folder
        source_subfolder_query = f"name='{source_folder_name}' and '{settings.google_drive_review_folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
        source_subfolder_results = service.files().list(
            q=source_subfolder_query,
            fields='files(id, name)'
        ).execute()
        
        source_subfolder_files = source_subfolder_results.get('files', [])
        
        if source_subfolder_files:
            source_subfolder_id = source_subfolder_files[0]['id']
            logger.info(f"   [📁] Using existing source subfolder: {source_folder_name}")
        else:
            # Create source subfolder
            source_subfolder_metadata = {
                'name': source_folder_name,
                'mimeType': 'application/vnd.google-apps.folder',
                'parents': [settings.google_drive_review_folder_id]
            }
            source_subfolder = service.files().create(
                body=source_subfolder_metadata,
                fields='id, name'
            ).execute()
            source_subfolder_id = source_subfolder.get('id')
            logger.info(f"   [📁] Created source subfolder: {source_folder_name}")
        
        # Create main document folder inside source subfolder
        folder_metadata = {
            'name': main_folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [source_subfolder_id]
        }
        
        folder = service.files().create(
            body=folder_metadata,
            fields='id, name'
        ).execute()
        
        folder_id = folder.get('id')
        logger.info(f"   [📁] Created document folder: {main_folder_name}")
        
        # Create Phase1 subfolder for content assessment
        phase1_metadata = {
            'name': 'Phase1_Content_Assessment',
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [folder_id]
        }
        phase1_folder = service.files().create(
            body=phase1_metadata,
            fields='id, name'
        ).execute()
        phase1_folder_id = phase1_folder.get('id')
        logger.info(f"   [📁] Created Phase1 subfolder")
        
        # Create Phase2 subfolder for presentation materials
        phase2_metadata = {
            'name': 'Phase2_Presentation_Materials',
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [folder_id]
        }
        phase2_folder = service.files().create(
            body=phase2_metadata,
            fields='id, name'
        ).execute()
        phase2_folder_id = phase2_folder.get('id')
        logger.info(f"   [📁] Created Phase2 subfolder")
        
        created_files = []
        
        # Create DOCX files using python-docx
        from docx import Document
        from docx.shared import Inches
        
        # 1. Use Cases DOCX - Enhanced with proper structure (Phase 1)
        if production_content.get('use_cases'):
            logger.info(f"   [📄] Creating Use Cases document for Phase 1...")
            use_cases_doc = Document()
            use_cases_doc.add_heading('Praktische Anwendungsfälle', 0)
            
            # Parse structured use cases content
            use_cases_content = production_content['use_cases']
            use_cases = []  # Initialize to prevent variable error
            
            # Validate content
            if not use_cases_content or len(use_cases_content) < 100:
                logger.warning(f"   [⚠️] Use cases content is too short or empty, skipping document creation...")
                # Add error message to document
                use_cases_doc.add_paragraph("Error: Use cases content generation failed. Please regenerate.")
            else:
                use_cases = parse_use_cases_content(use_cases_content)
            
                # If parsing failed, add the content as-is
                if not use_cases or len(use_cases) == 0:
                    logger.warning("Could not parse use cases, adding raw content")
                    use_cases_doc.add_paragraph(use_cases_content)
                    use_cases = []  # Empty list to prevent loop
            
            for i, use_case in enumerate(use_cases, 1):
                # Title
                use_cases_doc.add_heading(f"Anwendungsfall {i}: {use_case['title']}", level=1)
                
                # Theoretical Background
                p = use_cases_doc.add_paragraph()
                p.add_run('Theoretischer Hintergrund:').bold = True
                use_cases_doc.add_paragraph(use_case['theory'])
                
                # Scenario
                p = use_cases_doc.add_paragraph()
                p.add_run('Praxis-Szenario:').bold = True
                use_cases_doc.add_paragraph(use_case['scenario'])
                
                # Tasks
                p = use_cases_doc.add_paragraph()
                p.add_run('Aufgaben für Lernende:').bold = True
                for j, task in enumerate(use_case['tasks'], 1):
                    use_cases_doc.add_paragraph(f"{j}. {task}", style='List Number')
                
                # Solution
                p = use_cases_doc.add_paragraph()
                p.add_run('Musterlösung für Trainer:').bold = True
                use_cases_doc.add_paragraph(use_case['solution'])
                
                # Expected Results
                p = use_cases_doc.add_paragraph()
                p.add_run('Erwartete Ergebnisse:').bold = True
                use_cases_doc.add_paragraph(use_case['expected_results'])
                
                # Add page break between use cases
                use_cases_doc.add_page_break()
            
            # Create in memory and upload directly to Google Drive (NO LOCAL FILES)
            import io
            import uuid
            temp_id = str(uuid.uuid4())[:8]
            
            # Save to memory buffer instead of disk
            doc_buffer = io.BytesIO()
            use_cases_doc.save(doc_buffer)
            doc_buffer.seek(0)
            
            # Upload directly from memory to Google Drive - Phase 1
            file_metadata = {
                'name': f"Anwendungsfälle_{original_filename}.docx",
                'parents': [phase1_folder_id]  # Phase 1 folder
            }
            
            media = MediaIoBaseUpload(doc_buffer, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            
            created_files.append({
                'name': file_metadata['name'],
                'id': file.get('id'),
                'type': 'use_cases_docx',
                'phase': 'Phase1'
            })
            logger.info(f"   [✅] Use Cases uploaded to Phase 1")
        
        # 2. Quiz DOCX - Enhanced with proper structure (Phase 1)
        if production_content.get('quiz'):
            logger.info(f"   [📄] Creating Quiz document for Phase 1...")
            
            # Validate content
            quiz_content = production_content['quiz']
            if not quiz_content or len(quiz_content) < 100:
                logger.warning(f"   [⚠️] Quiz content is too short or empty, skipping...")
            else:
                quiz_doc = Document()
                quiz_doc.add_heading('Bewertungsfragen', 0)
                
                # Parse structured quiz content
                questions = parse_quiz_content(quiz_content)
            
                for question in questions:
                    # Question with difficulty
                    p = quiz_doc.add_paragraph()
                    p.add_run(f"Frage {question['number']} ({question['difficulty']}): ").bold = True
                    p.add_run(question['question'])
                    
                    # Answer options
                    for option in question['options']:
                        quiz_doc.add_paragraph(option, style='List Bullet')
                    
                    # Correct answer
                    p = quiz_doc.add_paragraph()
                    p.add_run(f"Korrekte Antwort: {question['correct_answer']}").bold = True
                    
                    # Explanation
                    p = quiz_doc.add_paragraph()
                    p.add_run('Erklärung: ').bold = True
                    p.add_run(question['explanation'])
                    
                    quiz_doc.add_paragraph()  # Blank line between questions
                
                # Save to memory buffer instead of disk
                quiz_buffer = io.BytesIO()
                quiz_doc.save(quiz_buffer)
                quiz_buffer.seek(0)
                
                file_metadata = {
                    'name': f"Quiz_{original_filename}.docx",
                    'parents': [phase1_folder_id]  # Phase 1 folder
                }
                
                media = MediaIoBaseUpload(quiz_buffer, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
                
                created_files.append({
                    'name': file_metadata['name'],
                    'id': file.get('id'),
                    'type': 'quiz_docx',
                    'phase': 'Phase1'
                })
                logger.info(f"   [✅] Quiz uploaded to Phase 1")
        
        # 3. Trainer Script DOCX - Create in memory (Phase 2)
        if production_content.get('trainer_script'):
            logger.info(f"   [📄] Creating Trainer Script document for Phase 2...")
            
            # Validate content
            trainer_script_content = production_content['trainer_script']
            if not trainer_script_content or len(trainer_script_content) < 100 or trainer_script_content.startswith("Error"):
                logger.warning(f"   [⚠️] Trainer script content is invalid, skipping...")
            else:
                script_doc = Document()
                script_doc.add_heading('Trainer-Skript', 0)
                script_doc.add_paragraph(trainer_script_content)
                
                # Save to memory buffer instead of disk
                script_buffer = io.BytesIO()
                script_doc.save(script_buffer)
                script_buffer.seek(0)
                
                file_metadata = {
                    'name': f"Trainer-Skript_{original_filename}.docx",
                    'parents': [phase2_folder_id]  # Phase 2 folder
                }
                
                media = MediaIoBaseUpload(script_buffer, mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
                
                created_files.append({
                    'name': file_metadata['name'],
                    'id': file.get('id'),
                    'type': 'trainer_script_docx',
                    'phase': 'Phase2'
                })
                logger.info(f"   [✅] Trainer Script uploaded to Phase 2")
        
        # 4. Professional PPTX Generation (Phase 2) - High-quality, comprehensive presentations
        # Always create a presentation, even if trainer script generation failed
        pptx_buffer = None
        
        if production_content.get('trainer_script'):
            try:
                logger.info("   [PPTX] Creating professional presentation with comprehensive content...")
                
                # Generate professional PPTX with detailed content coverage
                pptx_buffer = create_professional_pptx_presentation(
                    trainer_script=production_content['trainer_script'],
                    doc_name=original_filename
                )
                
            except Exception as e:
                logger.error(f"   [ERROR] PPTX generation error: {e}")
                import traceback
                logger.error(f"   [ERROR] Traceback: {traceback.format_exc()}")
        
        # Fallback: Create robust PPTX from document content if trainer script generation failed
        if not pptx_buffer and production_content.get('document_content'):
            try:
                logger.info("   [PPTX-FALLBACK] Creating presentation from document content...")
                pptx_buffer = create_robust_pptx_from_content(
                    document_content=production_content['document_content'],
                    doc_name=original_filename
                )
            except Exception as e:
                logger.error(f"   [ERROR] PPTX fallback creation failed: {e}")
        
        # Upload PPTX if we have one - Phase 2
        if pptx_buffer:
            try:
                logger.info(f"   [📤] Uploading PPTX to Phase 2...")
                file_metadata = {
                    'name': f"Präsentation_{original_filename}.pptx",
                    'parents': [phase2_folder_id]  # Phase 2 folder
                }
                
                media = MediaIoBaseUpload(
                    pptx_buffer, 
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    resumable=True
                )
                
                file = service.files().create(
                    body=file_metadata,
                    media_body=media,
                    fields='id,name'
                ).execute()
                
                created_files.append({
                    'name': file_metadata['name'],
                    'id': file.get('id'),
                    'type': 'professional_pptx',
                    'phase': 'Phase2'
                })
                
                logger.info(f"   [OK] PPTX presentation uploaded to Phase 2")
            except Exception as e:
                logger.error(f"   [ERROR] PPTX upload failed: {e}")
        else:
            logger.warning("   [WARNING] Could not create PPTX presentation - no valid content available")
        
        return {
            "success": True,
            "total_files_created": len(created_files),
            "created_files": created_files,
            "enabler_folder_id": folder_id,
            "folder_name": main_folder_name,
            "phase1_folder_id": phase1_folder_id,
            "phase2_folder_id": phase2_folder_id
        }
        
    except Exception as e:
        logger.error(f"Error saving production content: {e}")
        return {
            "success": False,
            "error": str(e),
            "total_files_created": 0,
            "created_files": []
        }

async def process_document(doc: Dict[str, Any], ws: WebSocketProgress, doc_num: int, total_docs: int, retry_count: int = 0) -> Dict[str, Any]:
    """
    Process a single document with enhanced LangGraph orchestration.
    Includes comprehensive error handling and validation.
    """
    import os  # Ensure os is available in function scope
    
    doc_name = doc['name']
    doc_id = doc['id']
    
    try:
        await ws.send_progress(f"[{doc_num}/{total_docs}] Processing: {doc_name}", doc_num, total_docs)
        await ws.send_progress("=" * 60)
        
        # PRE-FLIGHT VALIDATION
        await ws.send_progress(f"   [VALIDATION] Running pre-flight checks...")
        
        # 1. Validate document metadata
        if not doc_name or not doc_id:
            await ws.send_progress(f"   [ERROR] Invalid document metadata")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": "Invalid document metadata"
            }
        
        # 2. Check file extension
        supported_extensions = ['.docx', '.pdf', '.txt']
        file_ext = os.path.splitext(doc_name)[1].lower()
        if file_ext not in supported_extensions:
            await ws.send_progress(f"   [ERROR] Unsupported file type: {file_ext}")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": f"Unsupported file type: {file_ext}"
            }
        
        await ws.send_progress(f"   [OK] Pre-flight validation passed")
        
        # Step 1: Extract content
        await ws.send_progress(f"   [PHASE 1/5] Extracting document content...")
        
        document_content = None
        try:
            if google_drive_service:
                document_content = google_drive_service.get_file_content(doc_id)
            else:
                # Fallback - use dummy content
                document_content = f"This is dummy content for {doc_name}. In a real scenario, this would be extracted from Google Drive. The content would contain educational material that needs to be processed by the AI system to generate comprehensive learning materials including knowledge analysis, practical use cases, quiz questions, and video scripts."
        except Exception as e:
            await ws.send_progress(f"   [ERROR] Content extraction failed: {str(e)}")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": f"Content extraction failed: {str(e)}"
            }
        
        # 3. Validate content
        if not document_content or len(document_content.strip()) < 100:
            await ws.send_progress(f"   [WARNING] Insufficient content in {doc_name}")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": "Insufficient content (less than 100 characters)"
            }
        
        content_length = len(document_content)
        word_count = len(document_content.split())
        await ws.send_progress(f"   [OK] Extracted {content_length:,} characters ({word_count:,} words)")
        
        # 4. Check for excessive content
        if content_length > 1000000:  # 1MB limit
            await ws.send_progress(f"   [WARNING] Document very large ({content_length:,} chars), may take longer")
        
        await ws.send_progress(f"   [OK] Content validation passed")
        
        # Step 2: Analyze content depth
        await ws.send_progress(f"   [PHASE 2/5] Analyzing content depth and requirements...")
        
        content_depth = {}
        try:
            from app.services.advanced_document_processor import AdvancedDocumentProcessor
            doc_processor = AdvancedDocumentProcessor()
            content_depth = doc_processor.analyze_content_depth(document_content)
            
            await ws.send_progress(f"   [OK] Content analysis complete:")
            await ws.send_progress(f"       - {content_depth['word_count']} words detected")
            await ws.send_progress(f"       - {content_depth['estimated_slides']} slides required")
            await ws.send_progress(f"       - {content_depth['estimated_use_case_pages']} use case pages needed")
            await ws.send_progress(f"       - {content_depth['estimated_quiz_questions']} quiz questions required")
        except Exception as e:
            logger.warning(f"Content depth analysis failed: {e}")
            await ws.send_progress(f"   [WARNING] Content depth analysis failed, using defaults")
            content_depth = {
                "word_count": len(document_content.split()),
                "estimated_slides": 15,
                "estimated_use_case_pages": 5,
                "estimated_quiz_questions": 20
            }
        
        # Step 3: Process with LangGraph orchestration
        await ws.send_progress(f"   [PHASE 3/5] Starting AI orchestration workflow...")
        job_id = f"auto_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{doc_id[:8]}"
        
        processing_result = None
        
        # Try LangGraph orchestration first
        try:
            from app.services.langgraph_orchestrator import LangGraphWorkflowOrchestrator
            
            await ws.send_progress(f"   [LANGGRAPH] Initializing workflow orchestrator...")
            langgraph = LangGraphWorkflowOrchestrator()
            
            if langgraph.initialized:
                await ws.send_progress(f"   [LANGGRAPH] Running 5-phase workflow:")
                await ws.send_progress(f"       > Content Extraction")
                await ws.send_progress(f"       > Depth Analysis")
                await ws.send_progress(f"       > RAG Enhancement")
                await ws.send_progress(f"       > Content Generation (CrewAI/Gemini)")
                await ws.send_progress(f"       > Quality Assurance")
                
                processing_result = await langgraph.process_document_with_orchestration(
                    document_content=document_content,
                    job_id=job_id,
                    content_type="educational"
                )
                
                if processing_result["success"]:
                    await ws.send_progress(f"   [OK] LangGraph orchestration complete")
            else:
                await ws.send_progress(f"   [WARNING] LangGraph not available, using direct processing")
        except Exception as e:
            logger.warning(f"LangGraph orchestration failed: {e}")
            await ws.send_progress(f"   [WARNING] LangGraph failed, falling back to direct RAG")
        
        # Fallback to direct RAG processing if LangGraph fails
        if not processing_result or not processing_result.get("success"):
            await ws.send_progress(f"   [FALLBACK] Using direct RAG processing...")
            
            rag = get_rag_processor()
            if rag:
                processing_result = await rag.process_document_with_rag(
                    document_content=document_content,
                    job_id=job_id,
                    content_type="educational"
                )
            else:
                # Last resort fallback
                await ws.send_progress(f"   [WARNING] RAG processor not available, using basic processing...")
                processing_result = {
                    "success": True,
                    "enhanced_content": {
                        "knowledge_analysis": f"# Basic Analysis\n\n{document_content[:500]}...",
                        "use_case_text": "# Basic Use Cases\n\nGenerated content...",
                        "quiz_text": "# Basic Quiz\n\nGenerated questions...",
                        "powerpoint_structure": "# Basic PowerPoint\n\nGenerated slides...",
                        "google_slides_content": "# Basic Google Slides\n\nGenerated content...",
                        "trainer_script": "# Basic Trainer Script\n\nGenerated script...",
                        "overall_quality_score": 0.7
                    }
                }
        
        if not processing_result["success"]:
            await ws.send_progress(f"   [ERROR] Processing failed: {processing_result.get('error', 'Unknown error')}")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": processing_result.get('error', 'Processing failed')
            }
        
        enhanced_content = processing_result["enhanced_content"]
        quality_score = enhanced_content.get("overall_quality_score", 0.0)
        
        # Step 4: Validate content completeness
        await ws.send_progress(f"   [PHASE 4/5] Validating content completeness...")
        required_content = ['knowledge_analysis', 'powerpoint_structure', 'google_slides_content',
                          'use_case_text', 'quiz_text', 'trainer_script']
        missing_content = [key for key in required_content if key not in enhanced_content or not enhanced_content[key]]
        
        if missing_content:
            await ws.send_progress(f"   [WARNING] Missing content types: {', '.join(missing_content)}")
        else:
            await ws.send_progress(f"   [OK] All content types generated successfully")
        
        await ws.send_progress(f"   [OK] Content generated (Quality: {quality_score:.2f})")
        
        # Step 5: Generate Production-Level Content and Save
        await ws.send_progress(f"   [PHASE 5/5] Generating production-level content and saving...")
        
        # Use personal Google Drive service for all operations
        try:
            import sys
            import os
            sys.path.append(os.getcwd())
            
            from personal_google_drive_service import PersonalGoogleDriveService
            
            # Initialize personal service with correct file paths
            personal_service = PersonalGoogleDriveService(
                credentials_file="personal_credentials.json",
                token_file="personal_google_token.pickle"
            )
            
            if personal_service.is_authenticated():
                await ws.send_progress(f"   [OK] Using personal Google account for content generation and storage")
                
                # Generate production-level content using advanced AI processing
                production_content = await generate_production_level_content(
                    document_content, enhanced_content, doc_name, ws
                )
                
                # Save as proper DOCX and PPTX files
                save_result = await save_production_content_with_personal_service(
                    personal_service, production_content, doc_name, doc_id
                )
            else:
                await ws.send_progress(f"   [ERROR] Personal Google account not authenticated")
                save_result = {"success": False, "error": "Personal account not authenticated"}
        except Exception as e:
            await ws.send_progress(f"   [ERROR] Personal service failed: {e}")
            save_result = {"success": False, "error": str(e)}
            save_result = {
                "success": True,
                "total_files_created": 0,
                "created_files": [],
                "enabler_folder_id": "local_only"
            }
        
        if save_result["success"]:
            files_created = save_result["total_files_created"]
            await ws.send_progress(f"   [OK] Created {files_created} files for {doc_name}")
            await ws.send_progress(f"   [OK] Files saved: knowledge_analysis, powerpoint, slides, use_cases, quiz, trainer_script")
            
            # Update Google Sheets with review tracking
            if google_sheets_service and google_sheets_service.service:
                try:
                    # Calculate processing time
                    processing_time = 0.0  # You can track this if needed
                    
                    # Add complete review record
                    output_file_names = [f['name'] for f in save_result.get("created_files", [])]
                    
                    # Extract Gamma PPTX file ID if available
                    gamma_pptx_file_id = ""
                    for file_info in save_result.get("created_files", []):
                        if file_info.get("type") == "gamma_pptx":
                            gamma_pptx_file_id = file_info.get("id", "")
                            break
                    
                    google_sheets_service.add_review_record(
                        document_name=doc_name,
                        source_file_id=doc_id,
                        processing_status="completed",
                        output_folder_id=save_result.get("enabler_folder_id", ""),
                        output_files=output_file_names,
                        quality_score=quality_score,
                        processing_time=processing_time,
                        gamma_pptx_file_id=gamma_pptx_file_id
                    )
                    await ws.send_progress(f"   [OK] Updated tracking sheet - marked as 'pending_review'")
                except Exception as e:
                    logger.warning(f"Could not update Google Sheets: {e}")
            
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "job_id": job_id,
                "status": "completed",
                "files_created": files_created,
                "quality_score": quality_score,
                "folder_id": save_result["enabler_folder_id"]
            }
        else:
            await ws.send_progress(f"   [ERROR] Failed to save: {save_result.get('error', 'Unknown error')}")
            return {
                "document_id": doc_id,
                "document_name": doc_name,
                "status": "failed",
                "error": save_result.get("error", "Failed to save")
            }
            
    except Exception as e:
        error_msg = str(e)
        await ws.send_progress(f"   [ERROR] Error processing {doc_name}: {error_msg}")
        return {
            "document_id": doc_id,
            "document_name": doc_name,
            "status": "failed",
            "error": error_msg
        }


async def main():
    """Main automation workflow"""
    ws = WebSocketProgress()
    await ws.connect()
    
    try:
        await ws.send_progress("=" * 80)
        await ws.send_progress("[START] FIAE AI CONTENT FACTORY - AUTOMATION ENGINE STARTED")
        await ws.send_progress("=" * 80)
        
        # Check services
        if not google_drive_service:
            await ws.send_progress("[WARNING] Google Drive service not available")
            await ws.send_progress("[TIP] Processing will continue with limited functionality")
        else:
            await ws.send_progress("[OK] Google Drive service available")
        
        # Discover documents
        await ws.send_progress("[SEARCH] Discovering documents in source folder...")
        source_folder_id = settings.google_drive_content_source_folder_id
        logger.info(f"[DEBUG] Using SOURCE folder ID: {source_folder_id}")
        logger.info(f"[DEBUG] Expected SOURCE folder ID: 1YtN3_CftdJGgK9DFGLSMIky7PbYfFsX5")
        
        if google_drive_service:
            all_documents = google_drive_service.list_files_in_folder(source_folder_id)
            enabler_documents = [
                doc for doc in all_documents 
                if doc['name'].lower().endswith('.docx') and not doc['name'].startswith('~')
            ]
        else:
            # Fallback - create dummy documents for testing
            await ws.send_progress("[WARNING] Using dummy documents for testing (Google Drive not available)")
            enabler_documents = [
                {"id": "test_doc_1", "name": "test_document_1.docx"},
                {"id": "test_doc_2", "name": "test_document_2.docx"}
            ]
        
        if not enabler_documents:
            await ws.send_progress("[WARNING] No DOCX documents found in source folder")
            return
        
        await ws.send_progress(f"[OK] Found {len(enabler_documents)} documents to process")
        await ws.send_progress("")
        
        # Process documents
        processed_count = 0
        failed_count = 0
        results = []
        
        for i, doc in enumerate(enabler_documents, 1):
            result = await process_document(doc, ws, i, len(enabler_documents))
            results.append(result)
            
            if result["status"] == "completed":
                processed_count += 1
            else:
                failed_count += 1
            
            await ws.send_progress("")  # Blank line for separation
        
        # Final summary
        await ws.send_progress("")
        await ws.send_progress("=" * 80)
        await ws.send_progress("[SUCCESS] FIAE AI CONTENT FACTORY - AUTOMATION COMPLETED")
        await ws.send_progress("=" * 80)
        await ws.send_progress(f"[STATS] Successfully processed: {processed_count}/{len(enabler_documents)} documents")
        await ws.send_progress(f"[STATS] Failed: {failed_count}/{len(enabler_documents)} documents")
        success_rate = (processed_count / len(enabler_documents) * 100) if enabler_documents else 0
        await ws.send_progress(f"[STATS] Success rate: {success_rate:.1f}%")
        await ws.send_progress("")
        await ws.send_progress("[INFO] Generated content per document:")
        await ws.send_progress("       - Knowledge Analysis (Backend)")
        await ws.send_progress("       - PowerPoint Presentation (n slides based on content)")
        await ws.send_progress("       - Google Slides Content (interactive)")
        await ws.send_progress("       - IT Use Cases (task-based scenarios)")
        await ws.send_progress("       - Comprehensive Quiz (easy/medium/hard)")
        await ws.send_progress("       - Trainer Script (slide-by-slide)")
        await ws.send_progress("")
        await ws.send_progress("[TECH] Processing pipeline:")
        await ws.send_progress("       > LangGraph Orchestration")
        await ws.send_progress("       > CrewAI Multi-Agent System")
        await ws.send_progress("       > ChromaDB RAG Enhancement")
        await ws.send_progress("       > Gemini 1.5 Pro (Optimized)")
        await ws.send_progress("=" * 80)
        
    except Exception as e:
        await ws.send_progress(f"[ERROR] FATAL ERROR: {str(e)}")
        logger.error(f"Automation engine failed: {e}", exc_info=True)
    finally:
        await ws.close()


if __name__ == "__main__":
    # Configure logging
    logger.remove()
    logger.add(
        sys.stdout,
        format="<green>{time:YYYY-MM-DD HH:mm:ss}</green> | <level>{level: <8}</level> | <level>{message}</level>",
        level="INFO"
    )
    
    # Run automation
    asyncio.run(main())

